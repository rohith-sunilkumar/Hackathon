#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation about the Smart City Citizen Portal tools
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call(['pip3', 'install', 'python-pptx'])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def create_content_slide(prs, title, content_items):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def create_two_column_slide(prs, title, left_items, right_items):
    """Create a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Left column
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf_left = left.text_frame
    for item in left_items:
        p = tf_left.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    # Right column
    right = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    tf_right = right.text_frame
    for item in right_items:
        p = tf_right.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(prs, 
        "Smart City Citizen Portal",
        "Technology Stack & Development Tools\n\nAn AI-Powered Civic Engagement Platform")
    
    # Slide 2: Project Overview
    create_content_slide(prs, "Project Overview", [
        "Full-Stack Web Application for civic issue reporting",
        "AI-Powered automatic complaint categorization",
        "Real-Time Analytics dashboard for administrators",
        "Responsive Design for all devices",
        "",
        "📊 Statistics:",
        "• 30+ Files",
        "• 3,500+ Lines of Code",
        "• 10+ API Endpoints",
        "• 8 Major Features"
    ])
    
    # Slide 3: MERN Stack
    create_content_slide(prs, "MERN Stack Architecture", [
        "M - MongoDB: NoSQL Database",
        "E - Express.js: Web Framework",
        "R - React: Frontend Library",
        "N - Node.js: JavaScript Runtime",
        "",
        "✅ JavaScript Everywhere",
        "✅ Modern & Scalable",
        "✅ Large Community Support",
        "✅ Production Ready"
    ])
    
    # Slide 4: Frontend Technologies
    create_content_slide(prs, "Frontend Technologies", [
        "⚛️ React 18 - UI Library",
        "   • Component-based architecture",
        "   • Hooks & Context API",
        "   • Efficient rendering",
        "",
        "🎨 Tailwind CSS - Styling Framework",
        "   • Utility-first approach",
        "   • Responsive design",
        "   • Custom theme support",
        "",
        "🚀 Vite - Build Tool",
        "   • Lightning-fast HMR",
        "   • Optimized production builds"
    ])
    
    # Slide 5: Frontend Libraries
    create_content_slide(prs, "Frontend Libraries", [
        "🧭 React Router v6 - Client-side routing",
        "   • Protected routes",
        "   • Nested routing",
        "",
        "📡 Axios - HTTP Client",
        "   • Request/response interceptors",
        "   • Error handling",
        "",
        "📊 Recharts - Data Visualization",
        "   • Pie, Bar, Line charts",
        "   • Interactive tooltips",
        "",
        "🎯 Lucide React - Modern icon library"
    ])
    
    # Slide 6: Backend Technologies
    create_content_slide(prs, "Backend Technologies", [
        "🟢 Node.js v20 - JavaScript Runtime",
        "   • Fast & scalable",
        "   • Event-driven architecture",
        "   • ES Modules support",
        "",
        "⚡ Express.js - Web Framework",
        "   • Minimal & flexible",
        "   • Robust middleware support",
        "   • RESTful API design",
        "",
        "📦 npm - Package Manager",
        "   • 182 backend dependencies",
        "   • 462 frontend dependencies"
    ])
    
    # Slide 7: Database
    create_content_slide(prs, "Database Technology", [
        "🍃 MongoDB - NoSQL Database",
        "   • Flexible schema",
        "   • JSON-like documents",
        "   • Geospatial indexing",
        "   • Aggregation pipelines",
        "",
        "🔧 Mongoose ODM",
        "   • Schema validation",
        "   • Middleware hooks",
        "   • Query building",
        "   • Population (joins)",
        "",
        "📍 Geospatial Features",
        "   • 2dsphere indexing for location queries"
    ])
    
    # Slide 8: AI Integration
    create_content_slide(prs, "AI & Machine Learning", [
        "🤖 OpenAI API (GPT-4 Turbo)",
        "",
        "Capabilities:",
        "• Automatic categorization (6 categories)",
        "• Priority assignment (high/medium/low)",
        "• Confidence scoring (80-95% accuracy)",
        "• Reasoning explanation",
        "",
        "Categories:",
        "Road, Lighting, Waste, Safety, Water, Other",
        "",
        "✅ Fallback system if AI unavailable"
    ])
    
    # Slide 9: Authentication & Security
    create_content_slide(prs, "Authentication & Security", [
        "🔐 JSON Web Tokens (JWT)",
        "   • Stateless authentication",
        "   • Token expiration handling",
        "",
        "🔒 Bcrypt.js",
        "   • Password hashing",
        "   • Salt rounds for security",
        "",
        "🛡️ Security Features:",
        "   • CORS configuration",
        "   • Input validation (express-validator)",
        "   • Role-based access control",
        "   • XSS protection"
    ])
    
    # Slide 10: File Upload
    create_content_slide(prs, "File Upload & Storage", [
        "📤 Multer - File Upload Middleware",
        "   • Multi-file upload (up to 5 images)",
        "   • File type validation",
        "   • Size limits (5MB per file)",
        "   • Disk storage",
        "",
        "☁️ Cloudinary (Optional)",
        "   • Cloud image storage",
        "   • Image optimization",
        "   • CDN delivery",
        "",
        "💾 Local Storage Alternative",
        "   • Uploads directory",
        "   • Static file serving"
    ])
    
    # Slide 11: Development Tools
    create_content_slide(prs, "Development Tools", [
        "🔄 Nodemon - Auto-restart server",
        "   • Watch for file changes",
        "   • Faster development workflow",
        "",
        "📝 dotenv - Environment variables",
        "   • Secure configuration",
        "   • Keep secrets out of code",
        "",
        "✅ ESLint - Code linting",
        "   • Code consistency",
        "   • Error prevention",
        "",
        "🎯 ES Modules - Modern JavaScript",
        "   • Better tree-shaking",
        "   • Cleaner imports"
    ])
    
    # Slide 12: API Architecture
    create_content_slide(prs, "RESTful API Design", [
        "📡 API Endpoints (10+):",
        "",
        "Authentication:",
        "• POST /api/auth/register",
        "• POST /api/auth/login",
        "• GET /api/auth/me",
        "",
        "Reports:",
        "• POST /api/reports (with AI classification)",
        "• GET /api/reports (with filters)",
        "• GET /api/reports/:id",
        "• PATCH /api/reports/:id (admin)",
        "• DELETE /api/reports/:id (admin)",
        "",
        "Analytics:",
        "• GET /api/stats (admin dashboard)"
    ])
    
    # Slide 13: State Management
    create_content_slide(prs, "State Management", [
        "⚛️ React Context API",
        "   • Global state management",
        "   • AuthContext for user state",
        "   • No external library needed",
        "",
        "🔄 React Hooks",
        "   • useState - Component state",
        "   • useEffect - Side effects",
        "   • useNavigate - Routing",
        "   • useParams - URL parameters",
        "",
        "Benefits:",
        "   • Simple API",
        "   • Built into React",
        "   • Lightweight solution"
    ])
    
    # Slide 14: Data Visualization
    create_content_slide(prs, "Data Visualization with Recharts", [
        "📊 Charts Implemented:",
        "",
        "1. Pie Chart - Category distribution",
        "   • Visual breakdown of report types",
        "",
        "2. Bar Chart - Status breakdown",
        "   • Open, In-Progress, Resolved counts",
        "",
        "3. Line Chart - 30-day trend",
        "   • Historical report patterns",
        "",
        "Features:",
        "• Responsive design",
        "• Interactive tooltips",
        "• Custom colors",
        "• Legend support"
    ])
    
    # Slide 15: Responsive Design
    create_content_slide(prs, "Responsive Design", [
        "📱 Mobile-First Approach",
        "",
        "Breakpoints:",
        "• Mobile: < 640px",
        "• Tablet: 640px - 1024px",
        "• Desktop: > 1024px",
        "",
        "Tailwind CSS Features:",
        "• Utility classes for responsiveness",
        "• Custom color palette",
        "• Consistent spacing scale",
        "• Typography system",
        "",
        "✅ Works on all devices",
        "✅ Touch-friendly interface"
    ])
    
    # Slide 16: Database Seeding
    create_content_slide(prs, "Database Seeding", [
        "🌱 Custom Seed Script",
        "",
        "Creates Demo Data:",
        "• 2 demo users (citizen + admin)",
        "• 8 sample reports",
        "• Various statuses and categories",
        "• Realistic civic complaints",
        "• Time-distributed data",
        "",
        "Demo Credentials:",
        "• Citizen: citizen@demo.com / password",
        "• Admin: admin@demo.com / password",
        "",
        "Command: npm run seed"
    ])
    
    # Slide 17: Middleware Stack
    create_content_slide(prs, "Express Middleware Stack", [
        "🔧 Middleware Layers:",
        "",
        "1. CORS - Cross-origin requests",
        "2. Body Parser - JSON parsing",
        "3. Static Files - Serve uploads",
        "4. Request Logger - Track requests",
        "5. Auth Middleware - JWT verification",
        "6. Validation - Input sanitization",
        "7. Error Handler - Global error handling",
        "",
        "Benefits:",
        "• Modular architecture",
        "• Reusable components",
        "• Clean separation of concerns"
    ])
    
    # Slide 18: Environment Configuration
    create_content_slide(prs, "Environment Configuration", [
        "⚙️ Backend (.env):",
        "• PORT=5000",
        "• MONGODB_URI=mongodb://...",
        "• JWT_SECRET=secret_key",
        "• OPENAI_API_KEY=sk-...",
        "• NODE_ENV=development",
        "",
        "🌐 Frontend (.env):",
        "• VITE_API_URL=http://localhost:5000/api",
        "",
        "🔒 Security:",
        "• .gitignore for sensitive files",
        "• .env.example templates",
        "• No secrets in code"
    ])
    
    # Slide 19: Deployment Tools
    create_content_slide(prs, "Deployment Platforms", [
        "🚀 Backend Hosting:",
        "• Render (recommended)",
        "• Railway",
        "• Heroku",
        "• AWS/Azure/GCP",
        "",
        "🌐 Frontend Hosting:",
        "• Vercel (recommended)",
        "• Netlify",
        "• GitHub Pages",
        "• Cloudflare Pages",
        "",
        "💾 Database:",
        "• MongoDB Atlas (cloud)",
        "• Local MongoDB (development)"
    ])
    
    # Slide 20: Documentation
    create_content_slide(prs, "Comprehensive Documentation", [
        "📚 8 Documentation Files:",
        "",
        "1. START_HERE.md - Navigation guide",
        "2. GETTING_STARTED.md - Quick start (10 min)",
        "3. README.md - Complete documentation",
        "4. SETUP.md - Detailed installation",
        "5. STRUCTURE.md - Code organization",
        "6. FEATURES.md - Feature list",
        "7. DEPLOYMENT.md - Production guide",
        "8. PROJECT_SUMMARY.md - Technical overview",
        "",
        "Total: 15,000+ words of documentation"
    ])
    
    # Slide 21: Key Features Summary
    create_content_slide(prs, "Key Features Implemented", [
        "✅ User Authentication (JWT)",
        "✅ AI-Powered Classification (OpenAI GPT-4)",
        "✅ Image Upload (up to 5 per report)",
        "✅ Location Tracking (GPS + manual)",
        "✅ Admin Dashboard with Analytics",
        "✅ Real-time Status Tracking",
        "✅ Responsive Design",
        "✅ Role-Based Access Control",
        "✅ Data Visualization (Charts)",
        "✅ RESTful API",
        "✅ Database Seeding",
        "✅ Comprehensive Documentation"
    ])
    
    # Slide 22: Why These Tools?
    create_content_slide(prs, "Tool Selection Criteria", [
        "✅ Modern & Popular",
        "   • Active development",
        "   • Large community",
        "   • Good documentation",
        "",
        "⚡ Performance",
        "   • Fast build times",
        "   • Efficient runtime",
        "   • Optimized output",
        "",
        "👨‍💻 Developer Experience",
        "   • Easy to learn",
        "   • Good tooling",
        "   • Clear APIs",
        "",
        "🚀 Production Ready",
        "   • Battle-tested",
        "   • Secure",
        "   • Scalable"
    ])
    
    # Slide 23: Learning Outcomes
    create_content_slide(prs, "Skills & Learning Outcomes", [
        "🎓 Full-Stack Development:",
        "• MERN stack mastery",
        "• RESTful API design",
        "• Database modeling",
        "",
        "🤖 AI Integration:",
        "• OpenAI API usage",
        "• Prompt engineering",
        "• Fallback strategies",
        "",
        "🎨 Modern Web Development:",
        "• React hooks & context",
        "• Tailwind CSS",
        "• ES6+ JavaScript",
        "",
        "🚀 DevOps Basics:",
        "• Environment configuration",
        "• Deployment strategies"
    ])
    
    # Slide 24: Project Statistics
    create_two_column_slide(prs, "Project Statistics",
        [
            "📊 Backend:",
            "• Files: 15",
            "• Lines: ~1,500",
            "• Functions: 30+",
            "• Routes: 10+",
            "• Models: 2",
            "",
            "📦 Dependencies:",
            "• Backend: 182 packages",
            "• Frontend: 462 packages"
        ],
        [
            "📊 Frontend:",
            "• Files: 15",
            "• Lines: ~2,000",
            "• Components: 10+",
            "• Pages: 8",
            "",
            "📚 Documentation:",
            "• Guides: 8",
            "• Words: 15,000+",
            "• Examples: 50+"
        ])
    
    # Slide 25: Development Timeline
    create_content_slide(prs, "Development Timeline", [
        "⏱️ Hackathon Breakdown:",
        "",
        "Day 1 (8 hours):",
        "• Backend setup & database models",
        "• API routes & AI integration",
        "",
        "Day 2 (8 hours):",
        "• Frontend setup & React components",
        "• Form handling & authentication",
        "",
        "Day 3 (8 hours):",
        "• Admin dashboard & analytics",
        "• Testing, debugging & documentation",
        "",
        "Total: ~24 hours of development"
    ])
    
    # Slide 26: Future Enhancements
    create_content_slide(prs, "Future Enhancements & Tools", [
        "🔮 Planned Additions:",
        "",
        "Real-Time Features:",
        "• Socket.io for notifications",
        "• WebSockets for chat",
        "",
        "Testing:",
        "• Jest for unit tests",
        "• Cypress for E2E tests",
        "",
        "Monitoring:",
        "• Sentry for error tracking",
        "• Google Analytics",
        "",
        "Maps:",
        "• Mapbox GL JS",
        "• Interactive map visualization"
    ])
    
    # Slide 27: Best Practices
    create_content_slide(prs, "Best Practices Implemented", [
        "✅ Code Quality:",
        "• Modular architecture",
        "• Separation of concerns",
        "• DRY principle",
        "• Clear naming conventions",
        "",
        "🔒 Security:",
        "• Environment variables",
        "• Input validation",
        "• Password hashing",
        "• JWT authentication",
        "",
        "📚 Documentation:",
        "• Code comments",
        "• README files",
        "• API documentation"
    ])
    
    # Slide 28: Conclusion
    create_content_slide(prs, "Technology Stack Summary", [
        "🎯 Complete Stack:",
        "",
        "Frontend:",
        "• React + Tailwind CSS + Vite",
        "",
        "Backend:",
        "• Node.js + Express + MongoDB",
        "",
        "AI:",
        "• OpenAI GPT-4",
        "",
        "Tools:",
        "• 20+ libraries & frameworks",
        "• Modern development practices",
        "• Production-ready architecture",
        "",
        "✅ Ready for deployment!"
    ])
    
    # Slide 29: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = "Thank You!"
    subtitle_shape.text = "Smart City Citizen Portal\n\nBuilding Better Cities with Modern Technology\n\n📧 Questions?\n🌐 GitHub: [Your Repository]\n📚 Documentation: 8 comprehensive guides"
    
    # Save presentation
    output_file = "Smart_City_Portal_Tools_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
